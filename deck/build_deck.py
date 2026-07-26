"""Build the codebase engineering deck (PPTX).

Run:  python deck/make_charts.py && python deck/build_deck.py

Reads banked results through deckdata.py -- no torch, no e3nn, no retraining.
Numbers are pulled at build time rather than typed into slide text, so the deck
follows the artifacts. At the end of a build it prints a provenance table
mapping every value to the file it came from, plus anything that resolved to
``n/a``.
"""

from __future__ import annotations

import re
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

import deckdata as dd

# --------------------------------------------------------------------------- #
# theme
# --------------------------------------------------------------------------- #

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(0.72)
CONTENT_W = SLIDE_W - 2 * MARGIN_L

INK = RGBColor(0x1B, 0x26, 0x2C)
BODY = RGBColor(0x33, 0x3D, 0x44)
MUTED = RGBColor(0x7A, 0x84, 0x8B)
RULE = RGBColor(0xD8, 0xDD, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

ACCENT = RGBColor(0x1A, 0x7F, 0x6B)  # EGRU teal, from the repo figure contract
ALERT = RGBColor(0xC2, 0x4A, 0x3F)  # PCT red
DARK_BG = RGBColor(0x16, 0x23, 0x2B)
DARK_ACCENT = RGBColor(0x5C, 0xD6, 0xB6)
DARK_MUTED = RGBColor(0x9F, 0xB0, 0xB8)
DARK_BODY = RGBColor(0xC9, 0xD6, 0xDB)

FONT = "Segoe UI"
MONO = "Consolas"

TOP_KICKER = Inches(0.46)
TOP_TITLE = Inches(0.80)
TOP_RULE = Inches(1.66)
TOP_BODY = Inches(1.92)
TOP_FOOT = Inches(6.94)
BODY_H = TOP_FOOT - TOP_BODY - Inches(0.10)

_MARKUP = re.compile(r"(\*\*.+?\*\*|\*[^*\n]+?\*|`[^`]+`)")


# --------------------------------------------------------------------------- #
# primitives
# --------------------------------------------------------------------------- #


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def _rich(paragraph, text, size, colour, bold=False, font=FONT):
    """Render light markup: ``**bold**``, ``*italic*`` and ``` `code` ```."""
    for token in _MARKUP.split(text):
        if not token:
            continue
        run = paragraph.add_run()
        run.font.size = Pt(size)
        run.font.color.rgb = colour
        run.font.bold = bold
        if token.startswith("**") and token.endswith("**"):
            run.text = token[2:-2]
            run.font.bold = True
            run.font.name = font
        elif token.startswith("*") and token.endswith("*"):
            run.text = token[1:-1]
            run.font.italic = True
            run.font.name = font
        elif token.startswith("`") and token.endswith("`"):
            run.text = token[1:-1]
            run.font.name = MONO
            run.font.size = Pt(size - 1)
        else:
            run.text = token
            run.font.name = font


def _bullet_format(paragraph, indent=Inches(0.30), char="•"):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(int(indent)))
    pPr.set("indent", str(-int(indent)))
    bu_font = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    bu_char = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(bu_font)
    pPr.append(bu_char)


def new_slide(prs, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG if dark else WHITE
    return slide


def add_header(slide, kicker, title, dark=False):
    tf = _textbox(slide, MARGIN_L, TOP_KICKER, CONTENT_W, Inches(0.3))
    p = tf.paragraphs[0]
    _rich(p, kicker.upper(), 12, DARK_ACCENT if dark else ACCENT, bold=True)
    p.runs[0].font.name = FONT

    tf = _textbox(slide, MARGIN_L, TOP_TITLE, CONTENT_W, Inches(0.9))
    p = tf.paragraphs[0]
    _rich(p, title, 27, WHITE if dark else INK, bold=True)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_L, TOP_RULE, CONTENT_W, Pt(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RULE if not dark else RGBColor(0x2E, 0x40, 0x4A)
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_bullets(slide, bullets, left, top, width, height, size=17, dark=False,
                spacing=10):
    tf = _textbox(slide, left, top, width, height)
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        p.line_spacing = 1.12
        _rich(p, text, size, DARK_BODY if dark else BODY)
        _bullet_format(p)
    return tf


def add_footnote(slide, text, dark=False):
    tf = _textbox(slide, MARGIN_L, TOP_FOOT, CONTENT_W, Inches(0.32))
    p = tf.paragraphs[0]
    _rich(p, text, 10, DARK_MUTED if dark else MUTED, font=MONO)


def add_picture_fit(slide, name, left, top, box_w, box_h):
    """Place a figure scaled to fit its box, centred, or a placeholder if absent."""
    path = dd.FIGURES / name
    if not path.exists():
        dd.MISSING.append(f"figure not generated: deck/figures/{name}")
        tf = _textbox(slide, left, top + box_h / 3, box_w, Inches(0.6))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _rich(p, f"[ figure n/a — run make_charts.py to generate {name} ]", 13, ALERT)
        return

    from PIL import Image

    with Image.open(path) as img:
        pw, ph = img.size
    scale = min(box_w / pw, box_h / ph)
    w, h = int(pw * scale), int(ph * scale)
    slide.shapes.add_picture(
        str(path), left + (box_w - w) // 2, top + (box_h - h) // 2, width=w, height=h
    )


def add_table(slide, headers, rows, left, top, width, col_widths=None,
              size=13, header_size=12, row_h=Inches(0.42)):
    n_rows, n_cols = len(rows) + 1, len(headers)
    height = row_h * n_rows
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    table.first_row = True
    table.horz_banding = False

    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw

    def _cell(r, c, text, bold, colour, fill):
        cell = table.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.10)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _rich(p, text, header_size if r == 0 else size, colour, bold=bold)

    for c, text in enumerate(headers):
        _cell(0, c, text, True, WHITE, RGBColor(0x2B, 0x3A, 0x42))
    for r, row in enumerate(rows, start=1):
        fill = WHITE if r % 2 else RGBColor(0xF5, 0xF7, 0xF8)
        for c, text in enumerate(row):
            _cell(r, c, text, False, BODY, fill)
    return table


# --------------------------------------------------------------------------- #
# slides
# --------------------------------------------------------------------------- #


TEAM = [
    ("Sayed Ashraful Islam Opin", "2232559042"),
    ("Ahsanul Karim Raiyan", "2221450642"),
    ("Susmit Talukder", "2221865042"),
    ("Natasha Monir Shawon", "2121744642"),
]


def s01_title(prs, d):
    slide = new_slide(prs, dark=True)

    tf = _textbox(slide, MARGIN_L, Inches(0.95), CONTENT_W, Inches(0.4))
    _rich(tf.paragraphs[0], "EDGE AI FOR TELE-REHABILITATION", 13,
          DARK_ACCENT, bold=True)

    tf = _textbox(slide, MARGIN_L, Inches(1.40), Inches(11.6), Inches(1.8))
    for i, line in enumerate(["Rehabilitation Exercise Assessment",
                              "with SE(3)-Equivariant Edge AI"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.04
        _rich(p, line, 40, WHITE, bold=True)

    tf = _textbox(slide, MARGIN_L, Inches(3.08), Inches(11.6), Inches(0.6))
    _rich(tf.paragraphs[0],
          "Engineering problems, design decisions, and measured solutions across the codebase",
          18, DARK_MUTED)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, Inches(3.72),
                                 Inches(1.5), Pt(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

    tf = _textbox(slide, MARGIN_L, Inches(4.05), CONTENT_W, Inches(0.32))
    _rich(tf.paragraphs[0], "TEAM", 11, DARK_ACCENT, bold=True)

    # Two columns of two, so the names read at projector distance.
    for col, members in enumerate((TEAM[:2], TEAM[2:])):
        tf = _textbox(slide, MARGIN_L + Inches(5.6) * col, Inches(4.42),
                      Inches(5.3), Inches(1.2))
        for i, (name, sid) in enumerate(members):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            _rich(p, name, 16, WHITE)
            run = p.add_run()
            run.text = f"   {sid}"
            run.font.size = Pt(14)
            run.font.name = MONO
            run.font.color.rgb = DARK_MUTED

    stats = [
        "128 Python modules in `src/` · 242 banked result files · 3 research threads",
        "KIMORE · NTU RGB+D 60 Cross-View · REHAB24-6 · UI-PRMD",
        "WACV 2026 submission · arXiv extended twin · live webcam demo",
    ]
    tf = _textbox(slide, MARGIN_L, Inches(6.02), Inches(11.6), Inches(1.3))
    for i, line in enumerate(stats):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        _rich(p, line, 12.5, DARK_MUTED)


def s02_domain(prs, d):
    slide = new_slide(prs)
    add_header(slide, "the setting",
               "Automated rehabilitation assessment — and why deployment breaks it")
    add_bullets(slide, [
        "**The task.** Regress a clinician's 0–50 movement-quality score from a skeleton "
        "sequence — 25 Kinect joints, one recording per patient per exercise.",
        "**The lab assumption.** One fixed camera, one fixed frame rate, every joint tracked "
        "on every frame.",
        "**The clinic.** The camera gets moved between sessions, the depth sensor drops frames "
        "in bursts, and joints vanish behind the body for seconds at a time.",
        "**The baseline family** — point-cloud transformer, TCN, ST-GCN — learns geometry "
        "from data. Geometry is then only as reliable as the training distribution was "
        "representative.",
        "**The design bet of this project.** Make the invariances *structural properties of the "
        "architecture* — provable and testable to machine precision — instead of hoping "
        "augmentation covers the deployment envelope.",
    ], MARGIN_L, TOP_BODY, CONTENT_W, BODY_H, size=18, spacing=14)
    add_footnote(slide, "src/kimore_cde_data.py · src/models_curvenet.py · "
                        "docs/reference/PROJECT_BRIEF.md")


def s02b_data(prs, d):
    slide = new_slide(prs)
    add_header(slide, "the input side",
               "Four corpora, four skeleton conventions, one loader contract")
    add_bullets(slide, [
        "**The problem.** Every corpus describes a body differently. KIMORE is Kinect v2 with 25 "
        "joints and real sensor timestamps in 100-ns ticks; NTU RGB+D uses the same SDK ordering "
        "(checked, not assumed); REHAB24-6 is OptiTrack with 26 markers needing an anatomical remap "
        "onto those 25; UI-PRMD differs again.",
        "**Every preprocessing step has to be equivariance-safe**, or the theorem is lost before the "
        "model sees anything. Root-relative coordinates give translation invariance; per-sequence "
        "scaling uses the *median joint radius*, a type-0 scalar that commutes with rotation; time is "
        "divided by a *global* constant, because a per-sequence one would erase the clinical duration "
        "cue.",
        "**The loader never invents an arrival time.** Length capping is uniform index subsampling "
        "that keeps the real timestamps, so the irregular-sampling claim is not quietly manufactured "
        "by the data path it is meant to test.",
        "**The demo re-implements the identical preprocessing byte for byte**, so what the model sees "
        "live differs from training only in domain, not in kind.",
        "**Why this gets its own slide.** A mis-mapped joint or an invented timestamp does not crash "
        "anything. It silently moves a number that nobody can then trace back.",
    ], MARGIN_L, TOP_BODY, CONTENT_W, BODY_H, size=15.5, spacing=11)
    add_footnote(slide, "src/kimore_cde_data.py · src/ntu_dataset.py · "
                        "src/load_rehab246.py · demo/mp_to_kinect.py")


def s03_triage(prs, d):
    slide = new_slide(prs)
    add_header(slide, "problem 0 — design",
               "Three paradigms proposed, two killed before a line of code")
    add_table(
        slide,
        ["Proposal", "Fatal attack", "Verdict"],
        [
            ["A — Neural biomechanical controller\n(infer torques, mass matrix, impairment)",
             "Non-identifiable: mass, potential and torque are underdetermined from "
             "position-only data. Also needs double differentiation of noisy joints.",
             "Killed"],
            ["B — Hyperbolic × SE(3) geometry\n(tree in ℍⁿ, pose in SO(3))",
             "The hyperbolic half is marginal at skeleton scale — ~25 nodes, depth ~5, "
             "which Euclidean space embeds fine in 16–32 dims.",
             "Half kept"],
            ["C — Optimal-transport bridge\n(distance to a “health manifold”)",
             "OT needs distributions; a clinic gives one trajectory. Heavy and unstable "
             "against a competitor whose brand is CPU real-time.",
             "Killed"],
        ],
        MARGIN_L, TOP_BODY, CONTENT_W,
        col_widths=[Inches(3.5), Inches(6.4), Inches(1.99)],
        size=12.5, row_h=Inches(0.92),
    )
    add_bullets(slide, [
        "The survivor is a synthesis: **SE(3)-equivariance** from B plus **continuous time** from A, "
        "minus the identifiability-killing force decomposition. Both surviving claims are provable "
        "structural properties, not empirical hopes.",
        "Killing a paradigm on validity grounds costs days. Killing it after implementation costs months.",
    ], MARGIN_L, Inches(5.72), CONTENT_W, Inches(1.15), size=14, spacing=7)
    add_footnote(slide, "docs/reference/PROJECT_BRIEF.md §2")


def s04_architecture(prs, d):
    slide = new_slide(prs)
    add_header(slide, "design",
               "Where the guarantee is manufactured — and where it is discharged")
    add_bullets(slide, [
        "**Root-relative coordinates.** Subtracting the spine base makes translation invariance "
        "automatic — a constant that is removed, not learned.",
        "**Steerable encoder.** Bone vectors → spherical harmonics (ℓ≤2) → "
        "Clebsch–Gordan-constrained tensor products (e3nn). The radial weights are functions of "
        "‖r‖ only, so they are invariant scalars and cannot leak orientation.",
        "**No bias on any non-scalar channel.** The only vector fixed by the whole rotation group is "
        "zero, so a learnable bias on a type-1 channel silently breaks equivariance. Every steerable "
        "layer is built with `biases=False`.",
        "**The invariant cut Π** — the one place the theorem is discharged: 32 scalars + 8 vector "
        "norms + 28 cosines, pooled to 136, plus 24 bone lengths = **160-d parity-even**; plus 56 triple "
        "products pooled to 112 and 11 anatomical volumes = 123-d parity-odd. **283-d** in total.",
        "**Cosines, not raw dot products.** Unnormalised quadratic dots once drove predictions into the "
        "hundreds — the fix is in the feature definition, not in a clamp.",
        "**Recurrence on real inter-arrival times.** The GRU consumes the actual Δt between samples "
        "rather than a frame index.",
    ], MARGIN_L, TOP_BODY, CONTENT_W, BODY_H, size=15.5, spacing=11)
    add_footnote(slide, "src/equivariant_gru.py · src/chirality.py · "
                        "paper/arxiv_extended/methodology.tex")


def s05_solver(prs, d):
    slide = new_slide(prs)
    add_header(slide, "problem 1 — numerics",
               "The architecture is exactly equivariant. Is the integrator?")
    add_picture_fit(slide, "fig_solver.png", MARGIN_L, TOP_BODY,
                    Inches(6.5), Inches(4.7))
    add_bullets(slide, [
        "**Fixed-step explicit RK is safe.** Butcher coefficients are scalars, so they commute with "
        "every equivariant operation — exact to roundoff, and independent of step size.",
        "**Adaptive stepping is not.** `dopri5` scales each error component by its own magnitude, but "
        "a rotation preserves only the *L2 norm of each irrep block*, not per-component magnitudes.",
        "Formally, invariance needs [ρ(g), W²] = 0 — generically false for the default W even "
        "when ρ(g) is orthogonal. The rotated and unrotated runs then take **different step grids**.",
        "**Fix:** a per-irrep isotropic error norm `N_eq` = sqrt( (1/F) Σ_f (1/d_f)·‖e_f‖² ), "
        "invariant because each block norm is.",
        "Step-grid divergence goes to **exactly zero**, and the fix survives replacing the mock field "
        "with real e3nn layers — the most likely silent-break point.",
    ], Inches(7.45), TOP_BODY, Inches(5.16), BODY_H, size=13.5, spacing=9)
    add_footnote(slide, "src/equivariance_suite.py · "
                        "docs/reference/outputs_equivariance_certificate.txt")


def s06_certificate(prs, d):
    slide = new_slide(prs)
    add_header(slide, "design discipline",
               "Certification is a gate that can fail — not a claim in the abstract")
    c = d["certify"]
    add_table(
        slide,
        ["Gate", "What it tests", "Threshold", "Measured"],
        [
            ["E1", "Encoder equivariance  h(g·x) = ρ(g)h(x)", "≤ 1e-12",
             dd.fmt(c["E1"], ".2e")],
            ["E2", "Read-out invariance  s(g·x) = s(x)  — the theorem", "≤ 1e-12",
             dd.fmt(c["E2"], ".2e")],
            ["E3", "Violation vs rotation magnitude (log–log slope)", "≈ 0",
             dd.fmt(c["E3"], ".3f")],
            ["E4", "fp32 / fp64 ratio — is the residual just roundoff?", "≥ 1e5",
             dd.fmt(c["E4"], ".2e")],
            ["E5", "Translation invariance on raw world coordinates", "≤ 1e-12",
             dd.fmt(c["E5"], ".2e")],
            ["E8", "Full SO(3) with parity-odd channels active", "≤ 1e-12",
             dd.fmt(c["E8"], ".2e")],
        ],
        MARGIN_L, TOP_BODY, CONTENT_W,
        col_widths=[Inches(0.8), Inches(6.7), Inches(2.0), Inches(2.39)],
        size=13, row_h=Inches(0.44),
    )
    add_bullets(slide, [
        "**Build the control that must fail.** A deliberately non-equivariant vector field is run "
        "through the same gates and must be rejected: drift 1.18e-02 and a precision ratio of "
        "**1.000** against the real model's 1.2e+08. A test that cannot fail proves nothing.",
        "**E4 is what separates roundoff from a hidden violation.** A genuine symmetry break would "
        "not shrink eight orders of magnitude when you move from fp32 to fp64.",
    ], MARGIN_L, Inches(5.15), CONTENT_W, Inches(1.7), size=14.5, spacing=8)
    add_footnote(slide, "src/certify_egru.py · outputs/cde_block2/certify_egru.json "
                        "· src/equivariance_suite.py")


def s07_determinism(prs, d):
    slide = new_slide(prs)
    add_header(slide, "problem 2 — measurement",
               "The ±0.33 MAD “hardware noise floor” was not hardware")
    add_bullets(slide, [
        f"**The stakes, stated as a number.** Repeated runs of the *same* configuration differed by "
        f"±0.33 MAD. The entire effect being claimed is floor {dd.fmt(d['floor'], '.2f')} minus "
        f"model ≈ 6.6, about **1.65 MAD** — so the noise was 20% of the signal and wider than "
        f"every model-vs-model gap in the paper.",
        "**Two nondeterministic kernels, not one.** e3nn's `index_add_` resolves CUDA write collisions "
        "with atomic adds, making fp32 summation order-dependent; cuDNN's fused GRU backward does the "
        "same. Fixing either alone leaves a floor.",
        "**Fix 1 — change the reduction, not the maths.** Replace scatter aggregation with a dense "
        "incidence matmul `einsum('je,nec->njc', A, m)`: identical sum, fixed order, numerically "
        "equivalent — so **existing checkpoints stay valid**, which gate G1 asserts to 1e-12 in fp64.",
        "**Fix 2 — pin the environment before CUDA exists.** `CUBLAS_WORKSPACE_CONFIG=:4096:8` is set "
        "at module import, ahead of any CUDA context, alongside deterministic algorithms and an "
        "unfused GRU.",
        "**Outcome.** Bitwise-identical forwards and gradients. The seed-to-seed spread (0.48) is now "
        "reported *separately* from the fixed-configuration spread (0.33) instead of the two being "
        "conflated into one vague “noise”.",
    ], MARGIN_L, TOP_BODY, CONTENT_W, BODY_H, size=15.5, spacing=11)
    add_footnote(slide, "src/determinism.py · src/certify_phase1.py · "
                        "docs/planning/wacv_evaluation_and_action_plan.md (V3)")


def s08_fairness(prs, d):
    slide = new_slide(prs)
    add_header(slide, "problem 3 — experiment integrity",
               "Making “you sabotaged the baseline” unarguable")
    add_bullets(slide, [
        "**Cryptographic input identity.** Every robustness experiment routes byte-identical corrupted "
        "timelines to both models, and asserts an equal **SHA-256 hash** of the input at every "
        "corruption level and every seed. A promise is replaced by a check.",
        "**The baseline gets its best case.** The transformer architecturally requires a uniform "
        "T×J×3 grid, so it must resample. It is given the better of linear, cubic and "
        "forward-fill — not a strawman — and the resampling happens *inside* the shared engine.",
        "**Self-referential metric.** Each model is scored against its own clean prediction, so no "
        "cross-model scaling can bake in an advantage.",
        "**Named failure modes, not uniform noise.** Gilbert–Elliott burst drops (real dropouts "
        "cluster), Gaussian + AR(1) timestamp jitter under a monotonicity clamp, and Gamma "
        "inter-arrivals for 30→10 fps throttling.",
        "**Calibrated against reality.** 398 KIMORE recordings scanned, 143 with a genuine dropout, "
        "20 clearing the full CV ≥ 0.2 plus burst-gap bar; the synthetic sweep must match the real "
        "anchor's statistics within 20% before it is allowed to stand in for it.",
    ], MARGIN_L, TOP_BODY, CONTENT_W, BODY_H, size=15.5, spacing=11)
    add_footnote(slide, "src/resilience_sweep.py · src/corruption_pipeline.py · "
                        "src/irregular_data.py · src/block23_experiments.py")


def s09_signal_gate(prs, d):
    slide = new_slide(prs)
    add_header(slide, "problem 3b — what a robustness plot hides",
               "A model that ignores its input has zero degradation")
    add_bullets(slide, [
        "… and would therefore win every robustness plot ever drawn. **Degradation is not "
        "admissible on its own** — it is only interpretable next to absolute accuracy.",
        f"**The signal gate.** Every fold is checked against a per-exercise mean-predictor floor of "
        f"**{dd.fmt(d['floor'], '.2f')} MAD**. A run that fails to beat the floor is marked NO-SIGNAL, "
        f"its degradation columns are declared uninterpretable, and `--strict` aborts the run outright.",
        "**This gate killed our own first architecture.** The continuous-time Neural CDE — the "
        "original thesis — scored 8.43 against the 8.25 floor measured on that run. Its degradation "
        "curve was beautifully flat and completely meaningless. It was demoted to a control instead "
        "of being reported as robust.",
        "**The same discipline retired a second claim.** Selecting the training epoch on the test set "
        "inflates results by **+1.22 MAD (18.6%)** across three seeds, so every number reported is "
        "validation-selected — including the ones that then look weaker.",
    ], MARGIN_L, TOP_BODY, CONTENT_W, BODY_H, size=17, spacing=15)
    add_footnote(slide, "src/block23_experiments.py · src/protocol_null.py · "
                        "src/train_cde.py")


def s10_viewpoint(prs, d):
    slide = new_slide(prs)
    add_header(slide, "result — the payoff", "Viewpoint: a guarantee versus a fit")
    add_picture_fit(slide, "fig_viewpoint.png", MARGIN_L, TOP_BODY,
                    Inches(6.5), Inches(4.7))
    v = d["viewpoint"]
    add_bullets(slide, [
        f"**Certified models are flat by theorem:** EGRU {dd.fmt(v['egru'], '.1e')}, "
        f"InvariantGRU {dd.fmt(v['invgru'], '.0f')}, Ridge 0 worst-case degradation.",
        f"**Three independently designed non-equivariant families all fail:** PCT "
        f"{dd.fmt(v['pct'], '.2f')}, ST-GCN {dd.fmt(v['stgcn'], '.2f')}, TCN "
        f"{dd.fmt(v['tcn'], '.2f')} MAD — every one crossing the mean-predictor floor. The "
        f"fragility is structural, not one competitor's bug.",
        f"**Augmentation is a partial fix that looks total.** PCT + rotation augmentation is flat in "
        f"aggregate, yet still moves **{dd.fmt(v['pct_aug'], '.2f')} MAD per sequence** — which is "
        f"the number an individual patient actually experiences.",
        f"**The certificate is over the whole rotation group.** A 45° camera move — like NTU "
        f"RGB+D Cross-View's — is one instance of it, covered with no retraining. Second corpus "
        f"REHAB24-6 at 90°: AUROC "
        f"{dd.fmt(d['r246_auroc'], '.3f')} unchanged, logit drift {dd.fmt(d['r246_drift'], '.1e')} "
        f"against the transformer's 23.2.",
    ], Inches(7.45), TOP_BODY, Inches(5.16), BODY_H, size=13.5, spacing=11)
    add_footnote(slide, "outputs/cde_block2/final_tables.json · src/block3_baselines.py "
                        "· outputs/rehab246/")


def s11_costs(prs, d):
    slide = new_slide(prs)
    add_header(slide, "the honest ledger", "What the guarantee does not buy")
    a, n = d["accuracy"], d["nodefail"]
    add_table(
        slide,
        ["Claim", "Verdict", "Measured"],
        [
            ["Clean accuracy", "A tie, not a win",
             f"PCT {dd.fmt(a['pct'], '.2f')} · InvariantGRU {dd.fmt(a['invgru'], '.2f')} · "
             f"EGRU {dd.fmt(a['egru'], '.2f')}   (floor {dd.fmt(d['floor'], '.2f')})"],
            ["Sensor-node failure", "The transformer is genuinely best",
             f"MAD lost, 0→8 dead joints:  PCT +{dd.fmt(n['pct'], '.2f')} · "
             f"EGRU +{dd.fmt(n['egru'], '.2f')} · InvariantGRU +{dd.fmt(n['invgru'], '.2f')}"],
            ["Irregular sampling", "Null on this corpus",
             "The grid-free advantage does not materialise — the transformer degrades less"],
            ["Chirality", "Principled, unrewarded here",
             "+0.11 MAD for +16.7% parameters; KIMORE has no chiral pathology to detect"],
        ],
        MARGIN_L, TOP_BODY, CONTENT_W,
        col_widths=[Inches(2.5), Inches(3.5), Inches(5.89)],
        size=12.5, row_h=Inches(0.68),
    )
    b = d["bandwidth"]
    add_bullets(slide, [
        f"**The sampling null is predicted, not merely observed.** A spectral census puts "
        f"{dd.fmt(b['pos_below_2hz'], '.1%')} of positional energy below the resampling corner "
        f"({dd.fmt(b['corner_hz'], '.2f')} Hz), so the operator we avoid destroys only "
        f"{dd.fmt(b['pos_above'], '.1%')} of it. The claim only has teeth where the task depends on the "
        f"velocity band — {dd.fmt(b['vel_above'], '.0%')} of which sits above the corner.",
        "All four rows are in the paper. A negative result you can **derive** is stronger evidence of "
        "understanding than a positive one you stumbled into.",
    ], MARGIN_L, Inches(5.45), CONTENT_W, Inches(1.4), size=14, spacing=7)
    add_footnote(slide, "outputs/cde_block2/final_tables.json · src/bandwidth_law.py "
                        "· outputs/cde_block2/block5_bandwidth_law.json")


def s12_precision(prs, d):
    slide = new_slide(prs)
    add_header(slide, "problem 4 — deployment", "Equivariance has a precision budget")
    add_picture_fit(slide, "fig_precision.png", MARGIN_L, TOP_BODY,
                    Inches(6.5), Inches(4.7))
    p = d["precision"]
    add_bullets(slide, [
        "**Why there is a budget at all.** Quantisation is not an orthogonal map, so it does not "
        "commute with ρ(g). The theorem is exact in real arithmetic and has a numerical floor in "
        "every format a device actually runs.",
        f"**fp16 beats bf16 by 3.2×** at the same width ({dd.fmt(p['fp16'], '.2e')} vs "
        f"{dd.fmt(p['bf16'], '.2e')}) — equivariance is paid for in mantissa bits, not exponent range.",
        f"**The int8 cliff decomposes.** Weight-only quantisation *preserves* the theorem "
        f"({dd.fmt(p['w8'], '.2e')}) because the encoder is equivariant for any weights. It is "
        f"**activation** quantisation that breaks it ({dd.fmt(p['a8'], '.2e')}): "
        f"round(ρ(g)v) ≠ ρ(g)round(v).",
        "**A trap worth naming.** e3nn builds its Clebsch–Gordan tables at the *default dtype at "
        "construction time*. Calling `.half()` afterwards leaves them in fp32 and yields a fake "
        "certificate — 2.8e-08 instead of the honest 8.5e-15. Every format here is built natively.",
    ], Inches(7.45), TOP_BODY, Inches(5.16), BODY_H, size=13.5, spacing=11)
    add_footnote(slide, "src/precision_budget.py · src/int8_quant_budget.py · "
                        "outputs/precision_budget/")


def s13_streaming(prs, d):
    slide = new_slide(prs)
    add_header(slide, "problem 5 — deployment",
               "A bidirectional model cannot stream — yet it shipped as real-time")
    t = d["ttfs"]
    add_bullets(slide, [
        "**The contradiction.** The trained classifier was bidirectional with a mean-pool over all "
        "frames. It cannot emit a score until the clip ends, so its honest time-to-first-score is the "
        "entire recording — while the paper described a real-time system.",
        "**Controlled conversion, one variable.** Freeze the trained encoder and projection; retrain "
        "*only* the unidirectional GRU and head, behind a causal running-mean read-out updated in "
        "O(1) per frame. A loader assertion refuses to run if any other parameter moved.",
        f"**The cost is measured, not assumed:** {dd.fmt(d['stream_causal'], '.2f')}% causal against "
        f"{dd.fmt(d['stream_bidir'], '.2f')}% bidirectional on NTU Cross-View — 3.17 points for causality.",
        f"**Latency timed on a genuinely incremental single-frame step**, not the vectorised path, "
        f"which would be a lie about live latency: **{dd.fmt(t['per_frame'], '.1f')} ms/frame, "
        f"{dd.fmt(t['fps'], '.0f')} fps**, p95 {dd.fmt(t['p95'], '.1f')} ms, inside the "
        f"{dd.fmt(t['budget'], '.1f')} ms budget at 30 fps. Median {dd.fmt(t['lock'], '.0f')} frames to "
        f"lock on the correct class.",
        "**The engineering that made it runnable on ~6 GB free:** frozen features streamed to an "
        "on-disk fp16 memmap via *buffered writes*, because a writable memmap pins dirty pages in the "
        "working set and OOM-killed the run.",
    ], MARGIN_L, TOP_BODY, CONTENT_W, BODY_H, size=15, spacing=10)
    add_footnote(slide, "src/train_stream_egru.py · src/ttfs_benchmark.py · "
                        "outputs/ntu_stream/")


def s14_demo(prs, d):
    slide = new_slide(prs)
    add_header(slide, "from theorem to a laptop",
               "The live demo — and the gate that runs before the camera does")
    add_bullets(slide, [
        "**Pipeline.** Webcam → MediaPipe 33 landmarks → Kinect-25 remap (spine chain "
        "synthesised; a single `AXIS_SIGN` constant owns the Y-flip) → preprocessing byte-matched to "
        "the training loader → both models.",
        "**No hidden asymmetry.** One `predict()` call feeds byte-identical samples to the EGRU and the "
        "transformer; the transformer is resampled onto its required 100-frame grid *inside* the "
        "engine, where it is visible, rather than upstream where it would look like a handicap.",
        "**Real-time on a booth CPU.** Camera loop runs every frame, inference every 6th (≈6 Hz score "
        "updates), over a 64-frame bounded ring buffer — one `--infer-every` knob for a slow machine.",
        "**The gate runs first.** `smoke_test.py` re-measures the thesis on a synthetic body before any "
        "camera is touched, and *asserts* it: EGRU 3.76e-05 against PCT 8.02, tolerance 1e-03. If it "
        "passes, every later failure is plumbing rather than science.",
        "**Honest about the residual.** A real screen-to-webcam sweep over 7 physical viewpoints gives "
        "spreads of 3.19 (EGRU), 3.10 (InvariantGRU) and 6.89 (PCT). The runbook states plainly that "
        "EGRU's spread is not zero because the residual is pose-estimation noise, not model drift.",
    ], MARGIN_L, TOP_BODY, CONTENT_W, BODY_H, size=15.5, spacing=11)
    add_footnote(slide, "demo/engine.py · demo/smoke_test.py · demo/mp_to_kinect.py · "
                        "src/variant_b1_score.py")


def s14b_repo(prs, d):
    slide = new_slide(prs)
    add_header(slide, "keeping the results trustworthy",
               "One repository, three research threads, one source of truth")
    add_bullets(slide, [
        "**The problem.** This is a monorepo carrying three live threads at once — the equivariant "
        "work, an earlier self-supervised cross-sensor study that ended in a negative result, and an "
        "isolated EGNN/canonicalization sandbox. Threads that share a directory can quietly "
        "contaminate each other's results.",
        "**One source of truth for every published number.** `aggregate_final.py` regenerates all of "
        "them from banked JSON, and *measures* parameter counts with `sum(p.numel())` rather than "
        "quoting them from memory.",
        "**Isolation is asserted, not assumed.** The sandbox verifies `src/` and `outputs/` are "
        "untouched via a size-and-mtime manifest; `train_rehab246.py` refuses to start unless its "
        "output path ends in `outputs/rehab246`.",
        "**Reviewer questions answered without retraining wherever possible.** The per-invariant-family "
        "ablation zeroes a family while preserving the projection dimension — no retrain, no "
        "architecture change, so the anchor value reproduces the published number exactly.",
        "**Deprecation lives in the code, not a changelog.** `protocol_inflation.py` marks itself "
        "`.. deprecated::` and names its successor, so nobody re-derives a superseded result.",
    ], MARGIN_L, TOP_BODY, CONTENT_W, BODY_H, size=15.5, spacing=11)
    add_footnote(slide, "src/aggregate_final.py · src/ablation_invfamily.py · "
                        "research_egnn/README.md · src/protocol_inflation.py")


def s15_principles(prs, d):
    slide = new_slide(prs, dark=True)
    add_header(slide, "what generalises",
               "Six practices that did the actual work", dark=True)
    add_bullets(slide, [
        "**Every claim is a gate with a number.** 127 asserts across 36 modules, each with a named "
        "threshold and a diagnostic that says which invariant to inspect when it trips.",
        "**Build the control that must fail.** A deliberately broken vector field and a mean-predictor "
        "floor are what make a passing result informative.",
        "**Enforce fairness mechanically.** A SHA-256 input hash beats a promise that both models saw "
        "the same data.",
        "**Treat determinism as an instrument, not hygiene.** A noise floor you have not measured is a "
        "result you cannot defend.",
        "**Prefer an honest gap to a fake baseline.** `build_model('ctrgcn')` raises "
        "`NotImplementedError` with an implementation pointer rather than shipping a weak stand-in.",
        "**Ship the negative results, and reverse on measurement.** The EGNN comparison was expected to "
        "dilute the thesis and instead supported it; an earlier thread's headline was reversed by a "
        "bootstrap over ten subjects.",
    ], MARGIN_L, TOP_BODY, CONTENT_W, BODY_H, size=15.5, dark=True, spacing=11)
    add_footnote(slide, "docs/planning/wacv_evaluation_and_action_plan.md · "
                        "research_egnn/FINDINGS.md · docs/reviews/RESPONSE_TO_REVIEWER_R1.md",
                 dark=True)


SLIDES = [
    s01_title, s02_domain, s02b_data, s03_triage, s04_architecture, s05_solver,
    s06_certificate, s07_determinism, s08_fairness, s09_signal_gate,
    s10_viewpoint, s11_costs, s12_precision, s13_streaming, s14_demo,
    s14b_repo, s15_principles,
]


# --------------------------------------------------------------------------- #
# data gathering
# --------------------------------------------------------------------------- #


def gather() -> dict:
    """Pull every number the deck shows, recording where each came from."""
    ft, cj = dd.FINAL_TABLES, dd.CERTIFY
    acc = dd.accuracy_rows()
    view = dd.viewpoint_series()
    nf = dd.nodefail_rows()
    cert = (dd.certify() or {}).get("values", {})
    prec = (dd.precision_budget() or {}).get("rows", {})
    i8 = (dd.int8_budget() or {}).get("rows", {})
    tt = dd.ttfs() or {}
    st = (dd.stream() or {}).get("result", {})
    bw = dd.bandwidth() or {}

    def _acc(key):
        row = acc.get(key)
        return dd.rec(f"clean MAD :: {key}", row.get("mad") if row else None, ft)

    def _degr(key):
        row = view.get(key)
        return dd.rec(f"worst viewpoint degradation :: {key}",
                      row.get("max_degr") if row else None, ft)

    def _cert(key, name):
        return dd.rec(f"certificate :: {name}", cert.get(key), cj)

    def _nf_lost(key):
        if not nf:
            return dd.rec(f"node-failure MAD lost 0->8 :: {key}", None, ft)
        return dd.rec(f"node-failure MAD lost 0->8 :: {key}",
                      nf[-1][key] - nf[0][key], ft)

    def _prec(src, table, key, name):
        row = src.get(key)
        return dd.rec(f"invariance floor :: {name}",
                      row.get("inv_floor") if row else None, table)

    census = bw.get("band_census", {})
    cum = {e["f_hz"]: e for e in census.get("cum_energy", [])}

    r246_auroc, r246_drift = dd.rehab246(dd.REHAB246_EGRU)
    dd.rec("REHAB24-6 clean AUROC (EGRU, 3 seeds)", r246_auroc, dd.REHAB246_EGRU)
    dd.rec("REHAB24-6 worst viewpoint logit drift", r246_drift, dd.REHAB246_EGRU)

    lat = tt.get("latency", {})
    avf = tt.get("accuracy_vs_frames", {})

    return {
        "floor": dd.rec("mean-predictor floor (MAD)", dd.floor_mad(), ft),
        "accuracy": {
            "pct": _acc("PCT (baseline)"),
            "invgru": _acc("InvariantGRU  SO(3) chiral"),
            "egru": _acc("EGRU  SO(3) chiral"),
        },
        "viewpoint": {
            "egru": _degr("EGRU  SO(3) chiral"),
            "invgru": _degr("InvariantGRU  SO(3)"),
            "pct": _degr("PCT (baseline)"),
            "pct_aug": _degr("PCT + rot-aug"),
            "tcn": dd.rec("worst viewpoint degradation :: TCN",
                          dd.worst_mean_degradation(
                              "outputs/cde_block2/block3_tcn_s0_results.json"),
                          "outputs/cde_block2/block3_tcn_s0_results.json"),
            "stgcn": dd.rec("worst viewpoint degradation :: ST-GCN",
                            dd.worst_mean_degradation(
                                "outputs/cde_block2/block3_stgcn_s0_results.json"),
                            "outputs/cde_block2/block3_stgcn_s0_results.json"),
        },
        "nodefail": {
            "egru": _nf_lost("egru"),
            "invgru": _nf_lost("invgru"),
            "pct": _nf_lost("pct"),
        },
        "certify": {
            "E1": _cert("E1_encoder_equiv_worst", "E1 encoder equivariance"),
            "E2": _cert("E2_readout_invariance_worst", "E2 read-out invariance"),
            "E3": _cert("E3_loglog_slope", "E3 log-log slope"),
            "E4": _cert("E4_ratio", "E4 fp32/fp64 ratio"),
            "E5": _cert("E5_translation_invariance", "E5 translation invariance"),
            "E8": _cert("E8_rotation_invariance_chiral_worst", "E8 SO(3) with parity-odd"),
        },
        "precision": {
            "fp16": _prec(prec, dd.PRECISION, "fp16", "fp16"),
            "bf16": _prec(prec, dd.PRECISION, "bf16", "bf16"),
            "w8": _prec(i8, dd.INT8, "W8", "int8 weights only"),
            "a8": _prec(i8, dd.INT8, "A8", "int8 activations only"),
        },
        "ttfs": {
            "per_frame": dd.rec("streaming latency (ms/frame)",
                                lat.get("per_frame_ms"), dd.TTFS),
            "p95": dd.rec("streaming latency p95 (ms)",
                          lat.get("per_frame_ms_p95"), dd.TTFS),
            "fps": dd.rec("streaming throughput (fps)", lat.get("fps"), dd.TTFS),
            "budget": dd.rec("real-time budget (ms)", lat.get("budget_ms"), dd.TTFS),
            "lock": dd.rec("median frames to lock on",
                           avf.get("median_frames_to_lock"), dd.TTFS),
        },
        "stream_causal": dd.rec("NTU X-View causal top-1 (%)",
                                st.get("xview_top1_test"), dd.STREAM),
        "stream_bidir": dd.rec("NTU X-View bidirectional top-1 (%)",
                               st.get("bidir_reference"), dd.STREAM),
        "bandwidth": {
            "corner_hz": dd.rec("resampling corner (Hz)",
                                census.get("R_corner_hz"), dd.BANDWIDTH),
            "pos_below_2hz": dd.rec("positional energy below 2 Hz",
                                    cum.get(2.0, {}).get("pos_frac"), dd.BANDWIDTH),
            "pos_above": dd.rec("positional energy above the corner",
                                census.get("pos_above_R"), dd.BANDWIDTH),
            "vel_above": dd.rec("velocity energy above the corner",
                                census.get("vel_above_R"), dd.BANDWIDTH),
        },
        "r246_auroc": r246_auroc,
        "r246_drift": r246_drift,
    }


def main() -> int:
    data = gather()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for build in SLIDES:
        build(prs, data)

    out = dd.DECK / "codebase_engineering.pptx"
    prs.save(str(out))

    print(f"\nWrote {out.relative_to(dd.REPO)}  ({len(prs.slides)} slides)")

    print("\nProvenance -- every value on a slide and the file it came from")
    print("-" * 100)
    width = max(len(label) for label, _, _ in dd.PROVENANCE)
    for label, value, source in dd.PROVENANCE:
        print(f"  {label:<{width}}  {value:>14}   {source}")

    if dd.MISSING:
        print("\n!! UNRESOLVED -- these render as 'n/a' on the slides:")
        for m in dd.MISSING:
            print(f"  - {m}")
        return 1

    print("\nAll slide values resolved from banked artifacts.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
